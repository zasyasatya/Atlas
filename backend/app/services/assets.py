"""Dataset / deck / artifact ingestion with lightweight introspection."""
from __future__ import annotations

import csv
import hashlib
import io
import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from sqlmodel import Session, desc, select

from app.core.config import settings
from app.domain.enums import AssetKind
from app.domain.models import Asset, User

_SUBDIR = {
    AssetKind.DATASET: "datasets",
    AssetKind.DECK: "decks",
    AssetKind.ARTIFACT: "artifacts",
    AssetKind.IMAGE: "datasets",
}


def _safe_name(name: str) -> str:
    keep = "".join(c if c.isalnum() or c in "._- " else "_" for c in name).strip()
    return keep.replace(" ", "_")[:120] or "file.bin"


def _inspect_csv(raw: bytes) -> dict[str, Any]:
    try:
        text = raw.decode("utf-8", errors="replace")
        sample = text[:200_000]
        reader = csv.reader(io.StringIO(sample))
        rows = list(reader)[:2000]
        if not rows:
            return {}
        header = rows[0]
        body = rows[1:]
        total_rows = max(text.count("\n") - 1, len(body))
        return {
            "columns": header[:60],
            "rows_preview": body[:12],
            "row_count": total_rows,
            "column_count": len(header),
        }
    except Exception:  # pragma: no cover - defensive
        return {}


def _inspect_pptx(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        slides = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts = [
                sh.text_frame.text.strip()
                for sh in slide.shapes
                if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
            ]
            slides.append({"index": idx, "title": texts[0] if texts else f"Slide {idx}",
                           "bullets": texts[1:6]})
        return {"slide_count": len(slides), "slides": slides[:40]}
    except Exception:
        return {}


def store_asset(
    session: Session,
    *,
    filename: str,
    content: bytes,
    kind: AssetKind,
    user: User | None,
    topic_id: int | None = None,
    lesson_id: int | None = None,
    title: str = "",
    description: str = "",
    stage: str = "raw",
    content_type: str = "",
) -> Asset:
    settings.ensure_dirs()
    safe = _safe_name(filename)
    stamp = datetime.now(timezone.utc).strftime("%Y%m%d%H%M%S")
    folder = settings.storage_dir / _SUBDIR[kind]
    folder.mkdir(parents=True, exist_ok=True)
    path = folder / f"{stamp}_{safe}"
    path.write_bytes(content)

    preview: dict[str, Any] = {}
    rows = cols = slides = None
    lower = safe.lower()
    if lower.endswith((".csv", ".tsv", ".txt")):
        preview = _inspect_csv(content)
        rows, cols = preview.get("row_count"), preview.get("column_count")
    elif lower.endswith((".xlsx", ".xls")):
        try:
            import openpyxl

            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            ws = wb.active
            data = [[("" if c is None else str(c)) for c in r] for r in ws.iter_rows(max_row=13, values_only=True)]
            if data:
                preview = {"columns": data[0][:60], "rows_preview": data[1:13],
                           "row_count": ws.max_row - 1, "column_count": ws.max_column}
                rows, cols = ws.max_row - 1, ws.max_column
            wb.close()
        except Exception:
            preview = {}
    elif lower.endswith(".pptx"):
        preview = _inspect_pptx(path)
        slides = preview.get("slide_count")

    prev_versions = session.exec(
        select(Asset).where(Asset.topic_id == topic_id, Asset.kind == kind, Asset.title == (title or safe))
    ).all()

    asset = Asset(
        topic_id=topic_id,
        lesson_id=lesson_id,
        kind=kind,
        title=title or safe,
        description=description,
        filename=safe,
        stored_path=str(path),
        content_type=content_type,
        size_bytes=len(content),
        checksum=hashlib.sha256(content).hexdigest()[:32],
        version=len(prev_versions) + 1,
        row_count=rows,
        column_count=cols,
        slide_count=slides,
        preview_json=json.dumps(preview, ensure_ascii=False, default=str),
        stage=stage,
        uploaded_by=user.id if user else None,
    )
    session.add(asset)
    session.commit()
    session.refresh(asset)
    return asset


def list_assets(
    session: Session, *, topic_id: int | None = None, kind: AssetKind | None = None
) -> list[Asset]:
    stmt = select(Asset).order_by(desc(Asset.created_at))
    if topic_id is not None:
        stmt = stmt.where(Asset.topic_id == topic_id)
    if kind is not None:
        stmt = stmt.where(Asset.kind == kind)
    return list(session.exec(stmt))
